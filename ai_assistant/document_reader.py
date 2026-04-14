import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image


def read_pdf(path):

    reader = PdfReader(path)

    text = ""

    for page in reader.pages:
        text += page.extract_text()

    return text


def read_docx(path):

    doc = Document(path)

    text = ""

    for p in doc.paragraphs:
        text += p.text

    return text


def read_pptx(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text

    return text


def read_image(path):

    img = Image.open(path)

    text = pytesseract.image_to_string(img)

    return text


def read_file(path):

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return read_pdf(path)

    elif ext == ".docx":
        return read_docx(path)

    elif ext == ".pptx":
        return read_pptx(path)

    elif ext in [".png", ".jpg", ".jpeg"]:
        return read_image(path)

    else:
        return "Unsupported file format"